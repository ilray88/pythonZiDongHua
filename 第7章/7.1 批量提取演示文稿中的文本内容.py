from pptx import Presentation
from docx import Document
word_file = Document()
file_path = 'F:\\代码文件\\第7章\\员工管理制度.pptx'
ppt = Presentation(file_path)
for i in ppt.slides:
    for j in i.shapes:
        if j.has_text_frame:
            text_frame = j.text_frame
            for paragraph in text_frame.paragraphs:
                word_file.add_paragraph(paragraph.text)
save_path = 'F:\\代码文件\\第7章\\员工管理制度.docx'
word_file.save(save_path)
